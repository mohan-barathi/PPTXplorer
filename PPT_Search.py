import os
import threading
import ctypes
from Tkinter import *
global top
global inputFrame, resultFrame, progressFrame
global search_string, root_path, subDirVar, dispVar, ToolDescription
global rootTestBox, searchStringBox, SubDirCheckBox, dispBoxLabel, actionButton, resultListBox
global abortFlag, NoOfFiles

	
def createFreshFrames():
	createInputFrame()
	createResultFrame()
	createProgressFrame()
	createMenu()
	
def createMenu():
	global top, ToolDescription
	#ToolDescription.encode('utf-8')
	menubar = Menu(top)
	menuCmd = Menu(menubar, tearoff = 0)
	ToolDescription = "Tool : PPTExplorer \n\n Version : 1.1 \n\n License : GPLv2 \n\n"
	ToolDescription = ToolDescription + "Developer : Mohan Barathi Selvanayakam \n"
	ToolDescription = ToolDescription + "Contact   : MohanBarathi.Selvanayakam@in.bosch.com \n\n "
	ToolDescription = ToolDescription + ("Known Issues : Files with unicode characters in their name\ncannot be opened directly using 'open' from the window.\n").encode('utf-8')
	ToolDescription = ToolDescription + "To-do : Add scroll bars to result list box \n\n"
	ToolDescription = ToolDescription + " The source code will be published upon release of beta version, once the tool is stable."
	menuCmd.add_command(label = "About..", command = displayMsgBox)#"showinfo","About the tool..",Description))
	menubar.add_cascade(label="Menu", menu=menuCmd)
	top.config(menu=menubar)
	
def displayMsgBox(boxtype = "showinfo",title = "PptxExplorer",content = None):
	import tkMessageBox
	global ToolDescription
	if content is None:
		content = ToolDescription
	if boxtype is "showinfo":
		return tkMessageBox.showinfo(title,content)
	elif boxtype is "showwarning":
		return tkMessageBox.showwarning(title,content)
	elif boxtype is "showerror":
		return tkMessageBox.showerror(title,content)
	elif boxtype is "askquestion":
		return tkMessageBox.askquestion(title,content)
	elif boxtype is "askokcancel":
		return tkMessageBox.askokcancel(title,content)
	elif boxtype is "askretrycancel":
		return tkMessageBox.askretrycancel(title,content)
	else :
		return "error"

def createInputFrame(lockstatus = "unlocked", buttonName = "Search"):
	global inputFrame, top
	global rootTestBox, searchStringBox, SubDirCheckBox, actionButton
	global search_string, root_path, subDirVar
	try:
		inputFrame.destroy()
	except:
		donothing = 1
	inputFrame = Frame(top, bg = "white smoke", height = 110, width = 400)
	
	#first text box
	searchStringBox = Text(inputFrame, selectbackground = "blue", height = "1", width = "30")
	searchStringBox.place(x= 120, y = 10)
	searchStringBox.insert(END,search_string)
	if lockstatus is not "unlocked":
		searchStringBox.config(state=DISABLED, fg = 'grey')
	var = StringVar()
	searchStringBoxLabel = Label(inputFrame, textvariable=var, bg = "white smoke")
	var.set("Search String :")
	searchStringBoxLabel.place(x= 10, y= 10)
	
	#second text box
	rootTestBox = Text(inputFrame, selectbackground = "blue", height = "1", width = "30")
	rootTestBox.place(x= 120, y = 50)
	rootTestBox.insert(END,root_path)
	if lockstatus is not "unlocked":
		rootTestBox.config(state=DISABLED, fg = 'grey')
	var = StringVar()
	rootBoxLabel = Label(inputFrame, textvariable=var, bg = "white smoke")
	var.set("Root Folder Path :")
	rootBoxLabel.place(x= 10, y= 50)
	
	#Check Box
	SubDirCheckBox = Checkbutton(inputFrame,variable = subDirVar, text = "Consider Sub Directories", onvalue = 1, \
								offvalue = 0, height = 1, width = 20, bg = "white smoke")
	if lockstatus is not "unlocked":
		SubDirCheckBox.config(state=DISABLED)
	SubDirCheckBox.place(x=10, y= 80)
	
	#Button
	actionButton = Button(inputFrame, text =buttonName, height = 1, width = 6, relief= RAISED )
	if buttonName is "Search":
		actionButton.config(command= onClickSearchActionButton)
	elif buttonName is "Abort":
		actionButton.config(command= onClickAbortActionButton)
	actionButton.place(x=300,y=80)
	inputFrame.pack(side = TOP)
	
def onClickSearchActionButton():
	global inputFrame
	global abortFlag, NoOfFiles
	global rootTestBox, searchStringBox, SubDirCheckBox, resultListBox
	global search_string, root_path, subDirVar
	resultListBox.delete(0,resultListBox.size())
	root_path = rootTestBox.get("1.0",'end-1c')
	search_string = searchStringBox.get("1.0", 'end-1c')
	if IsErrorPresent():
		return
	listPptxFiles(root_path,subDirVar.get())
	#ctypes to show and confirm no. of files
	if 'no' in displayMsgBox("askquestion",str(NoOfFiles)+"pptx files Found..!",str(NoOfFiles)+" files found in specified path..!\n Do you want to proceed with the search..!?"):
		return
	else:
		abortFlag = False
		t = threading.Thread(target=loopedSearch)
		t.start()
		createInputFrame('locked', "Abort")
	
def onClickAbortActionButton():
	global abortFlag, actionButton
	abortFlag = True
	createInputFrame("unlocked", "Search")

def IsErrorPresent():
	global root_path,search_string,dispVar

	try:
		if ('\n' in search_string) or (len(search_string) is 0) or ('\t' in search_string):
			#include a ctype showing the error
			displayMsgBox('showerror', "ERROR", "Invalid Search String \n Make sure the string does not \n have a 'new line' or 'tab'")
			dispVar.set("Invalid Search String..!")
			search_string = ""
			createInputFrame("unlocked", "Search")
			return True
		elif not os.path.isdir(root_path):
			#include a ctype showing the error
			displayMsgBox('showerror', "ERROR", "Invalid Path\nMake sure the path is valid, and input characters does not\ncontain 'new line' or 'tab'")
			dispVar.set("Invalid Search Path..!")
			root_path = ""
			createInputFrame("unlocked", "Search")
			return True
		else:
			search_string = search_string.lower()
			return False
	except:
		#include ctype "pray to almighty"
		displayMsgBox('showerror', "ERROR", "Some unhandled exception Occured.\n Please contact the developer")
		return False

	
def createProgressFrame(displayString = "slm5cob/RBEI/ECP2"):
	global progressFrame, dispBoxLabel, dispVar
	try:
		progressFrame.destroy()
	except:
		donothing = 1
	progressFrame = Frame(top, bg = "white smoke", height = 40, width = 400)
	dispVar = StringVar()
	dispBoxLabel = Label(progressFrame, textvariable=dispVar, bg = "white smoke")
	dispVar.set(displayString)
	dispBoxLabel.place(x= 20, y= 10)
	progressFrame.pack(side = BOTTOM)
	
def createResultFrame():
	global resultFrame, resultListBox
	try:
		resultFrame.destroy()
	except:
		donothing = 1
	resultFrame = Frame(top, bg = "white smoke", height = 220, width = 400)
	resultListBox = Listbox(resultFrame, selectmode = SINGLE,  height=10, width = 58, relief= SUNKEN)
	resultListBox.place(x= 20, y= 10)
	#verticalScrollbar = Scrollbar(resultListBox, command = resultListBox.yview, width = 20)
	#resultListBox.config(yscrollcommand = verticalScrollbar.set)
	#verticalScrollbar.pack(side=RIGHT,fill=Y)
	openPptButton = Button(resultFrame, command = onClickOpenButton, text ="Open", height = 1, width = 6, relief= RAISED )
	openPptButton.place(x=300, y=180)
	resultFrame.pack(side = BOTTOM)

def onClickOpenButton():
	global resultListBox, root_path
	import subprocess
	
	try:
		selectedStr = resultListBox.get(resultListBox.curselection()[0])
	except:
		return
	selectedStr = selectedStr[8:]
	junk, requiredPath = selectedStr.split(":")
	fullFilePath = (root_path.encode('utf-8') + requiredPath.encode('utf-8'))
	try:	
		os.remove("openPPTx.bat")
	except:
		donothing = 1
	f = open("openPPTx.bat",'w+')
	f.write("mode con: cols=20 lines=1\n")
	#f.write("echo %1\n")
	#f.write("pause\n")
	f.write("%1\n")
	f.close()
	subprocess.Popen(["openPPTx.bat", fullFilePath])

	
def loopedSearch():
	global root_path,search_string,abortFlag,NoOfFiles,subDirVar,dispVar, actionButton
	presentFileCount = 0
	from fnmatch import fnmatch
	pattern = "*.pptx"
	for path, subdirs, files in os.walk(root_path):
		for name in files:
			if fnmatch(name, pattern):
				if abortFlag == True:
					dispVar.set("Search Aborted : " + str(presentFileCount) + "\\" + str(NoOfFiles))
					return
				presentFileCount = presentFileCount+1
				dispVar.set("Search in Progress : " + str(presentFileCount) + "\\" + str(NoOfFiles))
				SearchForString(os.path.join(path, name),search_string)
		if not subDirVar.get():
			break

	dispVar.set("Search Complete : " + str(presentFileCount) + "\\" + str(NoOfFiles))
	actionButton.config(text = "New Search", width = 10)
			


def listPptxFiles(root,IsSubDirTaken = 0):
	from fnmatch import fnmatch
	global NoOfFiles
	NoOfFiles = 0
	pattern = "*.pptx"
	for path, subdirs, files in os.walk(root):
		for name in files:
			if fnmatch(name, pattern):
				#print os.path.join(path, name)
				#FileList.append(os.path.join(path, name))
				NoOfFiles = NoOfFiles+1
		if not IsSubDirTaken:
			break


def SearchForString(prs_path,search_string):
	global root_path, resultListBox
	try:
		from pptx import Presentation
		prs = Presentation(prs_path)
		search_string = search_string.lower()
		for iter in range (0,len(prs.slides)):
			slide_text = ''
			for shape in prs.slides[iter].shapes:
				if abortFlag == True:
					return
				if not shape.has_text_frame:
					continue
				for paragraph in shape.text_frame.paragraphs:
					slide_text = slide_text + " " + paragraph.text
						
			if search_string in slide_text.lower():
				resultListBox.insert(1, ("Slide : " + (str(iter+1)) + (" of :") + (prs_path[len(root_path):].encode('utf-8'))))
	except:
		t = threading.Thread(target=ctypeDisplay, args=["Cannot read\n\""+prs_path+"\"\nThe File might be corrupted"])
		t.start()

			

def ctypeDisplay(displayString,titleString = "Error"):
	MB_OK = 0x0
	MB_OKCXL = 0x01
	MB_YESNOCXL = 0x03
	MB_YESNO = 0x04
	MB_HELP = 0x4000
	ICON_EXLAIM=0x30
	ICON_INFO = 0x40
	ICON_STOP = 0x10
	result = ctypes.windll.user32.MessageBoxA(0, displayString.encode('utf-8'), titleString, MB_OK | ICON_EXLAIM)
	return result
	
def startGUI():
	global top, subDirVar, abortFlag, search_string, root_path
	top = Tk()
	top.title("Pptx_Explorer 1.1")
	top.resizable(width=False, height=False)
	subDirVar = IntVar()
	abortFlag = False
	search_string = ""
	root_path = ""
	createFreshFrames()
	top.mainloop()
	try:	
		os.remove("openPPTx.bat")
	except:
		donothing = 1
	os.system("taskkill /F /PID " + str(os.getpid()))
	
if __name__ == '__main__':
	startGUI()